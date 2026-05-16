"""Build the Mom's Saheli pitch deck as a .pptx file.

Run:  python scripts/build_deck.py
Out:  docs/pitch_deck.pptx (5 slides, 16:9, brand-matched)

Layout philosophy:
  * Strict zones per slide (eyebrow / title / body / footer) so nothing collides.
  * Rich-text runs for italic-accent words — no two boxes ever stack on the
    same baseline (this was the slide-1 overlap bug).
  * Slides 2 + 4 + 5 inherit their working layout from v1; slides 1 + 3 are
    rewritten with way more meat per Avinash's feedback.
  * Self-explanatory: someone reading the deck without me presenting still
    understands what we built, how it works, and why it's not just "an AI bot".

Brand tokens mirror frontend/tailwind.config.js exactly.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── Brand tokens ──────────────────────────────────────────────────────────
BRAND_500 = RGBColor(0xF5, 0x9E, 0x0B)
BRAND_600 = RGBColor(0xD9, 0x77, 0x06)
BRAND_700 = RGBColor(0xB4, 0x53, 0x09)
BRAND_100 = RGBColor(0xFE, 0xF3, 0xC7)
BRAND_50  = RGBColor(0xFF, 0xFB, 0xEB)
CREAM     = RGBColor(0xFD, 0xF8, 0xF0)
INK_950   = RGBColor(0x09, 0x09, 0x0B)
INK_900   = RGBColor(0x18, 0x18, 0x1B)
INK_700   = RGBColor(0x3F, 0x3F, 0x46)
INK_500   = RGBColor(0x71, 0x71, 0x7A)
INK_300   = RGBColor(0xD4, 0xD4, 0xD1)
INK_200   = RGBColor(0xE5, 0xE5, 0xE3)
INK_100   = RGBColor(0xF4, 0xF4, 0xF3)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS   = RGBColor(0x15, 0x80, 0x3D)
SUCCESS_BG= RGBColor(0xDC, 0xFC, 0xE7)
DANGER    = RGBColor(0xB9, 0x1C, 0x1C)
DANGER_BG = RGBColor(0xFE, 0xE2, 0xE2)

SERIF = "Georgia"
SANS  = "Helvetica Neue"
MONO  = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_PATH = Path(__file__).resolve().parent.parent / "docs" / "pitch_deck.pptx"


# ═════════════════════════════════════════════════════════════════════════
# Layout primitives
# ═════════════════════════════════════════════════════════════════════════
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs, bg: RGBColor = CREAM):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = bg
    bg_rect.line.fill.background()
    bg_rect.shadow.inherit = False
    return s


def _apply_run_style(run, *, font=SANS, size=14, bold=False, italic=False,
                     color=INK_900, letter_spacing=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))


def add_text(slide, x, y, w, h, text, *, font=SANS, size=14, bold=False,
             italic=False, color=INK_900, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, spacing=1.15, letter_spacing=None):
    """Single-style text box. Newlines become separate paragraphs."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE  # critical: stop frames from growing past h
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, part in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = part
        _apply_run_style(r, font=font, size=size, bold=bold, italic=italic,
                         color=color, letter_spacing=letter_spacing)
    return tb


def add_rich(slide, x, y, w, h, runs, *, font=SANS, size=14, bold=False,
             color=INK_900, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             spacing=1.15):
    """Mixed-style text in ONE paragraph (the fix for slide-1 overlap).

    Each run is a dict with at minimum {'text': ...}; styling keys override
    the box defaults. Use '\\n' inside a run's text to break a paragraph.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE  # critical: stop frames from growing past h
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    paragraphs: list[list[dict]] = [[]]
    for run in runs:
        chunks = run["text"].split("\n")
        for i, chunk in enumerate(chunks):
            if i > 0:
                paragraphs.append([])
            r = dict(run)
            r["text"] = chunk
            paragraphs[-1].append(r)

    for pi, para_runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for run_def in para_runs:
            r = p.add_run()
            r.text = run_def["text"]
            _apply_run_style(
                r,
                font=run_def.get("font", font),
                size=run_def.get("size", size),
                bold=run_def.get("bold", bold),
                italic=run_def.get("italic", False),
                color=run_def.get("color", color),
                letter_spacing=run_def.get("letter_spacing"),
            )
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            s.adjustments[0] = max(0.0, min(0.5, radius))
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_eyebrow(slide, x, y, text, color=BRAND_700):
    add_text(slide, x, y, 12, 0.3, text.upper(),
             font=SANS, size=10, bold=True, color=color, letter_spacing=2.0)


def add_footer(slide, idx, total=5, dark=False):
    fg = INK_500
    add_text(slide, 0.6, 7.15, 8, 0.3,
             "Mom's Saheli  ·  Agent Forge AI Hackathon  ·  San Francisco  ·  May 16 2026",
             font=SANS, size=9, color=fg)
    add_text(slide, 12.0, 7.15, 1.0, 0.3, f"{idx:02d} / {total:02d}",
             font=MONO, size=9, color=fg, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE (rebuilt: single-frame headline, useful right panel)
# ═════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    s = add_slide(prs, bg=CREAM)

    # Right panel — replace the empty washes with a real "What it is" card
    add_rect(s, 8.95, 0, 4.38, 7.5, fill=BRAND_50)

    # Top-left brand bar
    add_text(s, 0.7, 0.55, 0.8, 0.8, "🌸",
             font=SERIF, size=44, color=BRAND_500)
    add_eyebrow(s, 1.45, 0.85,
                "Agent Forge AI Hackathon  ·  San Francisco  ·  May 16 2026")

    # SINGLE-FRAME HEADLINE — italic accent is a run inside the same paragraph
    # so it can never overlap a sibling box. Manual line break placed where
    # both halves fit cleanly inside the 8.2" box at 40pt Georgia bold.
    add_rich(s, 0.7, 1.85, 8.2, 3.1,
             runs=[
                 {"text": "The friend every working\nmom can "},
                 {"text": "finally afford", "italic": True, "color": BRAND_700},
                 {"text": "."},
             ],
             font=SERIF, size=40, bold=True, color=INK_950, spacing=1.15)

    # Subhead — clear product sentence
    add_text(s, 0.7, 4.95, 8.2, 1.3,
             "A 5-agent AI swarm that delivers $2,000 of consultant work "
             "— market research, compliance review, copywriting, launch — "
             "free, for America's 7 million working single moms.",
             font=SANS, size=16, color=INK_700, spacing=1.4)

    # Status badges
    add_rect(s, 0.7, 6.4, 0.18, 0.18, fill=SUCCESS, radius=0.5)
    add_text(s, 0.95, 6.33, 7, 0.3,
             "LIVE DEMO  ·  REAL GEMINI  ·  REAL CITED .GOV LAW  ·  REAL LAUNCH PAGES",
             font=SANS, size=9, bold=True, color=INK_700, letter_spacing=1.5)

    # ─── Right panel content (the actually-useful version) ───
    add_eyebrow(s, 9.2, 0.85, "What it does", color=BRAND_700)

    benefits = [
        ("Profiles", "her skills, hours, gap, and constraints"),
        ("Finds", "6–10 realistic income paths from real listings"),
        ("Blocks", "anything illegal — cites the actual .gov law on screen"),
        ("Writes", "her offer, copy, price, and 7-day launch plan"),
        ("Publishes", "a real shareable landing page in 38 seconds"),
    ]
    by = 1.35
    for verb, rest in benefits:
        add_rich(s, 9.2, by, 3.9, 0.55,
                 runs=[
                     {"text": "✓  ", "color": SUCCESS, "bold": True, "size": 13},
                     {"text": verb + " ", "bold": True, "color": INK_900, "size": 13},
                     {"text": rest, "color": INK_700, "size": 13},
                 ],
                 font=SANS, spacing=1.3)
        by += 0.7

    # Author block bottom of right panel
    add_text(s, 9.2, 5.75, 3.9, 0.3, "BUILT IN 24 HOURS  ·  SOLO",
             font=SANS, size=9, bold=True, color=BRAND_700, letter_spacing=2.0)
    add_text(s, 9.2, 6.1, 3.9, 0.5, "Avinash Ahuja",
             font=SERIF, size=20, bold=True, color=INK_900)
    add_text(s, 9.2, 6.55, 3.9, 0.35, "github.com/Avinash07-git/momsaheli",
             font=MONO, size=10, color=INK_500)

    add_footer(s, 1)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM (keep working layout, tighten copy)
# ═════════════════════════════════════════════════════════════════════════
def slide_problem(prs):
    s = add_slide(prs, bg=WHITE)
    add_eyebrow(s, 0.7, 0.6, "The world's pain")

    add_rich(s, 0.7, 1.05, 12, 2.4,
             runs=[
                 {"text": "She's not short on effort.\n"},
                 {"text": "She's short on margin.", "italic": True, "color": BRAND_700},
             ],
             font=SERIF, size=44, bold=True, color=INK_950, spacing=1.05)

    # Sub-paragraph that EXPLAINS the framing
    add_text(s, 0.7, 3.3, 12, 0.5,
             "America's working single moms aren't undermotivated. They're under-margined. "
             "Every dollar of help they need is gated by an expert who charges more than they earn.",
             font=SANS, size=13, italic=True, color=INK_500, spacing=1.4)

    stats = [
        ("7.3M",  "single moms in\nAmerica today",                  "Center for American Progress"),
        ("75%",   "already working —\nmost full-time",              "CAP"),
        ("$40K",  "median income for a\nworking single mom",        "CAP"),
        ("35%",   "of that income eaten\nby childcare alone",       "Child Care Aware of America"),
    ]
    card_w, card_h, gap = 2.85, 2.2, 0.15
    y = 4.05
    for i, (num, label, src) in enumerate(stats):
        x = 0.7 + i * (card_w + gap)
        add_rect(s, x, y, card_w, card_h, fill=CREAM, line=INK_100, radius=0.05)
        add_text(s, x + 0.25, y + 0.2, card_w - 0.5, 1.0, num,
                 font=SERIF, size=44, bold=True, color=BRAND_700, spacing=0.95)
        add_text(s, x + 0.25, y + 1.1, card_w - 0.5, 0.7, label,
                 font=SANS, size=12, color=INK_900, spacing=1.25)
        add_text(s, x + 0.25, y + card_h - 0.4, card_w - 0.5, 0.25,
                 f"SOURCE: {src}".upper(),
                 font=SANS, size=7.5, bold=True, color=INK_500, letter_spacing=1.5)

    # Brutal arithmetic closer
    add_rect(s, 0.7, 6.5, 11.95, 0.5, fill=INK_950, radius=0.15)
    add_rich(s, 0.9, 6.54, 11.6, 0.45,
             runs=[
                 {"text": "Help she needs:  ", "color": INK_300, "size": 11},
                 {"text": "consultant $2,000  ·  lawyer $200/hr  ·  marketer $200/hr  ·  bookkeeper $200/hr",
                  "color": CREAM, "bold": True, "size": 11},
                 {"text": "        Help she can afford:  ", "color": INK_300, "size": 11},
                 {"text": "$0", "color": BRAND_500, "bold": True, "size": 14},
             ],
             font=SANS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 2)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION (major rewrite: meaty 5-card grid + system intro)
# ═════════════════════════════════════════════════════════════════════════
def slide_solution(prs):
    s = add_slide(prs, bg=CREAM)
    add_eyebrow(s, 0.7, 0.45, "Our solution  ·  Mom's Saheli")

    # Headline + subhead establish the SYSTEM, not just "AI"
    add_rich(s, 0.7, 0.85, 12, 1.4,
             runs=[
                 {"text": "5 specialist AI agents do the\n$2,000 of consultant work — "},
                 {"text": "free", "italic": True, "color": BRAND_700},
                 {"text": "."},
             ],
             font=SERIF, size=32, bold=True, color=INK_950, spacing=1.05)

    add_text(s, 0.7, 2.35, 12, 0.5,
             "Not a chatbot. A swarm of specialists. Each agent has one narrow job, "
             "uses real production tools (not mocks), and hands its output to the next. "
             "End-to-end in 38 seconds, in your browser.",
             font=SANS, size=13, italic=True, color=INK_500, spacing=1.4)

    # 5 agent cards — each role description capped to ~110 chars (2 lines @ 10pt)
    agents = [
        {
            "n": "1", "name": "Profile",
            "role": "Reads her persona — skills, hours, $/mo gap, hard "
                    "constraints — into a typed schema.",
            "tools": "Pydantic schemas",
            "out": "Profile object",
        },
        {
            "n": "2", "name": "Market Scout",
            "role": "Pulls 6–10 real evidence cards (Etsy, Castiron, FB groups) "
                    "and asks Gemini to rank by net $/month.",
            "tools": "Bright Data + Actionbook + Gemini",
            "out": "Ranked opportunities",
        },
        {
            "n": "3", "name": "Reality &\nCompliance",
            "role": "Web-searches the state's actual cottage-food law and "
                    "BLOCKS illegal options with the live .gov citation.",
            "tools": "Tavily / Bright Data + Gemini",
            "out": "PASS / BLOCK + cited law",
        },
        {
            "n": "4", "name": "Launch",
            "role": "Gemini writes her offer, copy, price, target customer, "
                    "and 7-day plan; publishes a real landing page.",
            "tools": "Gemini + Actionbook + Butterbase",
            "out": "Real /launch/{slug} URL",
        },
        {
            "n": "5", "name": "Memory",
            "role": "Persists her run, then surfaces a cross-user pattern \u2014 "
                    "\u201cmoms like you in CA pivoted to X (71% win).\u201d",
            "tools": "Evermind",
            "out": "Persisted run + pattern",
        },
    ]

    card_w = 2.36
    gap = 0.1
    card_y = 3.05
    card_h = 3.55
    for i, a in enumerate(agents):
        x = 0.7 + i * (card_w + gap)
        # Card surface
        add_rect(s, x, card_y, card_w, card_h, fill=WHITE, line=INK_100, radius=0.04)
        # Top brand strip (taller to fit 2-line names like "Reality &\nCompliance")
        add_rect(s, x, card_y, card_w, 0.7, fill=BRAND_700, radius=0.04)
        # Number badge
        add_text(s, x + 0.18, card_y + 0.08, 0.55, 0.55, a["n"],
                 font=SERIF, size=26, bold=True, color=CREAM,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Agent name (right of number)
        add_text(s, x + 0.78, card_y + 0.08, card_w - 0.88, 0.55, a["name"],
                 font=SANS, size=11.5, bold=True, color=CREAM,
                 spacing=1.05, anchor=MSO_ANCHOR.MIDDLE)
        # Role description — 2 lines max, anchored mid for vertical centering
        add_text(s, x + 0.2, card_y + 0.85, card_w - 0.4, 1.55, a["role"],
                 font=SANS, size=9.5, color=INK_700, spacing=1.35)
        # Divider
        add_rect(s, x + 0.2, card_y + 2.5, card_w - 0.4, 0.015, fill=INK_200)
        # Tools row
        add_text(s, x + 0.2, card_y + 2.6, card_w - 0.4, 0.2, "USES",
                 font=SANS, size=7.5, bold=True, color=INK_500, letter_spacing=1.5)
        add_text(s, x + 0.2, card_y + 2.78, card_w - 0.4, 0.3, a["tools"],
                 font=MONO, size=8, color=INK_900, spacing=1.15)
        # Output row
        add_text(s, x + 0.2, card_y + 3.07, card_w - 0.4, 0.2, "OUTPUTS",
                 font=SANS, size=7.5, bold=True, color=INK_500, letter_spacing=1.5)
        add_text(s, x + 0.2, card_y + 3.25, card_w - 0.4, 0.22, "→ " + a["out"],
                 font=SANS, size=8.5, bold=True, color=BRAND_700)

    # Bottom strip — the demo proof
    add_rect(s, 0.7, 6.78, 11.95, 0.32, fill=INK_950, radius=0.12)
    add_text(s, 0.7, 6.79, 11.95, 0.3,
             "Live at  localhost:5173/run/jenny  ·  38 seconds end-to-end  ·  every step streams to the UI via SSE",
             font=MONO, size=9.5, color=CREAM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 3)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE PROOF (keep working layout, sharpen copy)
# ═════════════════════════════════════════════════════════════════════════
def slide_proof(prs):
    s = add_slide(prs, bg=WHITE)
    add_eyebrow(s, 0.7, 0.6, "Proof  ·  nothing is hardcoded")

    add_rich(s, 0.7, 1.0, 12, 1.4,
             runs=[
                 {"text": "Two moms. Same engine.\n"},
                 {"text": "Completely different output.", "italic": True, "color": BRAND_700},
             ],
             font=SERIF, size=36, bold=True, color=INK_950, spacing=1.05)

    add_text(s, 0.7, 2.85, 12, 0.45,
             "Click Run Jenny or Run Jessica. Same 5 agents. Different state, "
             "different skills, different gap. The system produces different launches — "
             "because it's actually reasoning, not retrieving canned demos.",
             font=SANS, size=12, italic=True, color=INK_500, spacing=1.4)

    card_y = 3.6
    card_h = 3.15
    gap = 0.3
    card_w = (12.6 - gap) / 2
    jenny_x = 0.7
    jess_x = jenny_x + card_w + gap

    def persona_card(x, emoji, name, sub, color_accent, steps):
        add_rect(s, x, card_y, card_w, card_h, fill=CREAM, line=INK_200, radius=0.04)
        add_rect(s, x, card_y, card_w, 0.85, fill=color_accent, radius=0.04)
        add_text(s, x + 0.25, card_y + 0.18, 1.0, 0.5, emoji,
                 font=SERIF, size=28, color=WHITE)
        add_text(s, x + 1.1, card_y + 0.18, card_w - 1.2, 0.35, name,
                 font=SERIF, size=20, bold=True, color=WHITE)
        add_text(s, x + 1.1, card_y + 0.50, card_w - 1.2, 0.3, sub,
                 font=SANS, size=10, color=WHITE)
        step_y = card_y + 1.05
        for i, (icon, text) in enumerate(steps):
            add_text(s, x + 0.3, step_y + i * 0.40, 0.3, 0.35, icon,
                     font=SANS, size=11, color=INK_900)
            add_text(s, x + 0.65, step_y + i * 0.40, card_w - 0.85, 0.36, text,
                     font=SANS, size=10, color=INK_700, spacing=1.25)

    persona_card(
        jenny_x, "👩🏽‍🍳", "Jenny",
        "Daycare aide  ·  California  ·  $600/mo gap  ·  5 hr/wk",
        BRAND_700,
        [
            ("→", "Top opp: weekend tiffin / meal-pack subscription ($720/mo)"),
            ("⚖", "BLOCKED: CA Health & Safety §114365 (no cottage-food)"),
            ("↩", "Pivots to #2: Sunday freezer-meal drops at farmers' market"),
            ("🚀", "Real launch page published at /launch/jenny-tiffin-fix"),
            ("✓", "Projected net: $620/mo  →  closes her gap"),
        ],
    )
    persona_card(
        jess_x, "💻", "Jessica",
        "Customer-service rep (WFH)  ·  Texas  ·  $400/mo gap  ·  3 hr/wk",
        INK_900,
        [
            ("→", "Top opp: Etsy printable kids' lunchbox planner pack"),
            ("✓", "Zero compliance hits — digital, async, no permits"),
            ("✓", "Constraint math passes (3 hr/wk is feasible)"),
            ("🚀", "Real Etsy listing draft + landing page published"),
            ("✓", "Projected net: $430/mo  →  closes her gap"),
        ],
    )

    add_footer(s, 4)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE ASK (add a 'what we built' recap)
# ═════════════════════════════════════════════════════════════════════════
def slide_ask(prs):
    s = add_slide(prs, bg=INK_950)
    add_rect(s, 8.5, 0,  4.83, 7.5, fill=INK_900)
    add_rect(s, 11.5, 0, 1.83, 7.5, fill=BRAND_700)

    add_eyebrow(s, 0.7, 0.7, "The ask", color=BRAND_500)

    add_text(s, 0.7, 1.15, 11.5, 1.6, "Pick us.",
             font=SERIF, size=84, bold=True, color=CREAM, spacing=1.0)

    add_rich(s, 0.7, 2.95, 11.5, 1.4,
             runs=[
                 {"text": "We're not shipping AI features.\nWe're closing the "},
                 {"text": "margin gap", "italic": True, "color": BRAND_500},
                 {"text": "."},
             ],
             font=SERIF, size=28, color=CREAM, spacing=1.2)

    # "What we built" recap — concrete proof of work
    add_eyebrow(s, 0.7, 4.45, "What shipped in 24 hours", color=BRAND_500)
    add_rich(s, 0.7, 4.78, 8.0, 0.5,
             runs=[
                 {"text": "5-agent swarm  ·  ", "color": CREAM, "size": 12},
                 {"text": "real Gemini calls", "color": BRAND_500, "size": 12, "bold": True},
                 {"text": "  ·  ", "color": CREAM, "size": 12},
                 {"text": "real cited .gov law", "color": BRAND_500, "size": 12, "bold": True},
                 {"text": "  ·  real published landing pages  ·  ", "color": CREAM, "size": 12},
                 {"text": "AgentField nested waterfall", "color": BRAND_500, "size": 12, "bold": True},
                 {"text": "  ·  SSE-streamed UI", "color": CREAM, "size": 12},
             ],
             font=SANS, spacing=1.4)

    # Sponsor strip
    sponsors = ["AgentField", "Bright Data", "Actionbook", "Evermind", "Butterbase",
                "Qwen Cloud", "Z.ai", "TokenRouter", "Zeabur"]
    add_text(s, 0.7, 5.5, 11.5, 0.3, "BUILT ON THE AGENT FORGE STACK · EVERY TOOL LOAD-BEARING",
             font=SANS, size=10, bold=True, color=INK_500, letter_spacing=2.0)
    add_rect(s, 0.7, 5.85, 11.95, 0.55, fill=INK_900, line=INK_700, radius=0.15)
    add_text(s, 0.7, 5.92, 11.95, 0.45, "  ·  ".join(sponsors),
             font=SANS, size=11, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # CTA card
    add_rect(s, 0.7, 6.6, 7.5, 0.55, fill=BRAND_500, radius=0.1)
    add_rich(s, 0.9, 6.66, 7.1, 0.45,
             runs=[
                 {"text": "TRY IT NOW  ", "color": INK_950, "size": 10, "bold": True, "letter_spacing": 2.0},
                 {"text": "→  github.com/Avinash07-git/momsaheli",
                  "color": INK_950, "size": 13, "bold": True, "font": MONO},
             ],
             font=SANS, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 11.4, 6.4, 1.5, 1.0, "🌸",
             font=SERIF, size=54, color=CREAM, align=PP_ALIGN.CENTER)

    add_text(s, 0.6, 7.15, 12.0, 0.3,
             "🐶 Built with love for every mom doing the math at 11pm.",
             font=SANS, size=10, italic=True, color=INK_500)
    add_text(s, 12.0, 7.15, 1.0, 0.3, "05 / 05",
             font=MONO, size=9, color=INK_500, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════
def main():
    prs = new_deck()
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_proof(prs)
    slide_ask(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"✅ Saved {OUT_PATH}")
    print(f"   {OUT_PATH.stat().st_size / 1024:.1f} KB · {len(prs.slides)} slides · 16:9")


if __name__ == "__main__":
    main()
